"""Office, OpenDocument and rich-text previewers."""

from __future__ import annotations

import base64
import html
import re
import zipfile
from pathlib import Path

from PyQt6.QtWidgets import QLabel, QTextBrowser, QVBoxLayout

from vibe_viewer.viewers.base import BaseViewer, ViewerError
from vibe_viewer.viewers.helpers import read_text_safely


class OfficeDocumentViewer(BaseViewer):
    name = "Office documents"
    category = "Office and spreadsheets"
    priority = 82
    extensions = (".docx", ".docm", ".dotx", ".pptx", ".pptm", ".potx", ".ppsx", ".odt", ".odp", ".rtf", ".doc", ".ppt")

    def __init__(self, parent=None) -> None:
        super().__init__(parent)
        self.info = QLabel()
        self.browser = QTextBrowser()
        self.browser.setOpenExternalLinks(False)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(6, 6, 6, 6)
        layout.addWidget(self.info)
        layout.addWidget(self.browser)

    def load_file(self, path: Path) -> None:
        suffix = path.suffix.lower()
        try:
            if suffix in {".docx", ".docm", ".dotx"}:
                content, details = self._render_docx(path)
            elif suffix in {".pptx", ".pptm", ".potx", ".ppsx"}:
                content, details = self._render_pptx(path)
            elif suffix in {".odt", ".odp"}:
                content, details = self._render_odf(path)
            elif suffix == ".rtf":
                content, details = self._render_rtf(path)
            else:
                content, details = self._render_legacy_ole(path)
        except Exception as exc:
            raise ViewerError(f"Не удалось открыть офисный документ: {exc}") from exc
        self.info.setText(f"{path.name} • {details}")
        self.browser.setHtml(content)

    @staticmethod
    def _render_docx(path: Path) -> tuple[str, str]:
        from docx import Document

        document = Document(path)
        output = ["<style>body{font-family:sans-serif;line-height:1.45} table{border-collapse:collapse}td{border:1px solid #aaa;padding:5px}img{max-width:90%}</style>"]
        for paragraph in document.paragraphs:
            text = html.escape(paragraph.text)
            style = (paragraph.style.name if paragraph.style else "").lower()
            if style.startswith("heading"):
                level_match = re.search(r"(\d+)", style)
                level = min(int(level_match.group(1)), 6) if level_match else 2
                output.append(f"<h{level}>{text}</h{level}>")
            elif text:
                output.append(f"<p>{text}</p>")
        for table in document.tables:
            output.append("<table>")
            for row in table.rows:
                output.append("<tr>" + "".join(f"<td>{html.escape(cell.text)}</td>" for cell in row.cells) + "</tr>")
            output.append("</table>")

        image_count = 0
        with zipfile.ZipFile(path) as archive:
            for name in archive.namelist():
                if not name.startswith("word/media/"):
                    continue
                data = archive.read(name)
                mime = _image_mime(name)
                output.append(
                    f'<p><img src="data:{mime};base64,{base64.b64encode(data).decode()}"></p>'
                )
                image_count += 1
        return "".join(output), f"DOCX • {len(document.paragraphs)} абзацев • {len(document.tables)} таблиц • {image_count} изображений"

    @staticmethod
    def _render_pptx(path: Path) -> tuple[str, str]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(path)
        output = ["<style>body{font-family:sans-serif}.slide{border:1px solid #ccd4df;border-radius:8px;margin:10px;padding:14px;background:white}img{max-width:85%;max-height:360px}</style>"]
        image_count = 0
        for index, slide in enumerate(presentation.slides, start=1):
            output.append(f'<section class="slide"><h2>Слайд {index}</h2>')
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = html.escape(shape.text).replace("\n", "<br>")
                    if text:
                        output.append(f"<p>{text}</p>")
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    encoded = base64.b64encode(image.blob).decode()
                    output.append(f'<p><img src="data:{image.content_type};base64,{encoded}"></p>')
                    image_count += 1
            output.append("</section>")
        return "".join(output), f"PowerPoint • {len(presentation.slides)} слайдов • {image_count} изображений"

    @staticmethod
    def _render_odf(path: Path) -> tuple[str, str]:
        from odf import teletype
        from odf.opendocument import load
        from odf.text import H, P

        document = load(str(path))
        output: list[str] = ["<style>body{font-family:sans-serif;line-height:1.45}</style>"]
        count = 0
        for node in document.getElementsByType(H):
            output.append(f"<h2>{html.escape(teletype.extractText(node))}</h2>")
            count += 1
        for node in document.getElementsByType(P):
            text = teletype.extractText(node)
            if text.strip():
                output.append(f"<p>{html.escape(text)}</p>")
                count += 1
        return "".join(output), f"OpenDocument • {count} текстовых блоков"

    @staticmethod
    def _render_rtf(path: Path) -> tuple[str, str]:
        from striprtf.striprtf import rtf_to_text

        raw, encoding, truncated = read_text_safely(path)
        text = rtf_to_text(raw)
        note = " • файл обрезан" if truncated else ""
        return f"<pre>{html.escape(text)}</pre>", f"RTF • {encoding}{note}"

    @staticmethod
    def _render_legacy_ole(path: Path) -> tuple[str, str]:
        """Best-effort text extraction for legacy DOC/PPT without external apps."""
        import olefile

        with olefile.OleFileIO(path) as ole:
            streams = ole.listdir(streams=True, storages=False)
            chunks: list[bytes] = []
            preferred = {"worddocument", "powerpoint document", "text_content"}
            for stream_parts in streams:
                if stream_parts[-1].lower() in preferred:
                    chunks.append(ole.openstream(stream_parts).read())
            if not chunks:
                chunks = [ole.openstream(parts).read() for parts in streams[:8]]
        raw = b"\n".join(chunks)
        candidates: list[str] = []
        for encoding in ("utf-16-le", "cp1251", "latin-1"):
            decoded = raw.decode(encoding, errors="ignore")
            runs = re.findall(r"[\w\s.,:;!?()\-–—«»'\"/@№%+]{5,}", decoded, flags=re.UNICODE)
            text = "\n".join(line.strip() for line in runs if line.strip())
            candidates.append(text)
        extracted = max(candidates, key=len, default="")[:2_000_000]
        if not extracted:
            extracted = "Текстовые фрагменты не найдены. Доступна только экспериментальная поддержка старого бинарного формата."
        return f"<pre>{html.escape(extracted)}</pre>", "legacy OLE • экспериментальное извлечение текста"


def _image_mime(name: str) -> str:
    extension = Path(name).suffix.lower()
    return {
        ".png": "image/png",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".svg": "image/svg+xml",
        ".tif": "image/tiff",
        ".tiff": "image/tiff",
    }.get(extension, "image/jpeg")

