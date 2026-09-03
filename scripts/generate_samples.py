"""Generate safe demo files without downloading anything."""

from __future__ import annotations

import csv
import gzip
import json
import sqlite3
import struct
import sys
import tarfile
import wave
import zipfile
from email.message import EmailMessage
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "samples" / "generated"


def write_text_samples() -> None:
    (OUTPUT / "hello.txt").write_text("Привет, Vibe Viewer!\nUniversal file preview.\n", encoding="utf-8")
    (OUTPUT / "example.md").write_text(
        "# Markdown\n\n- встроенный просмотр\n- таблицы\n- исходный код\n\n```python\nprint('hello')\n```\n",
        encoding="utf-8",
    )
    (OUTPUT / "data.json").write_text(
        json.dumps({"project": "Vibe Viewer", "formats": ["JSON", "PDF", "DOCX"]}, ensure_ascii=False),
        encoding="utf-8",
    )
    (OUTPUT / "config.yaml").write_text("project: Vibe Viewer\nready: true\nscore: 100\n", encoding="utf-8")
    (OUTPUT / "page.html").write_text("<h1>HTML preview</h1><p>Работает внутри Qt.</p>", encoding="utf-8")
    with (OUTPUT / "table.csv").open("w", encoding="utf-8", newline="") as stream:
        writer = csv.writer(stream)
        writer.writerow(["Формат", "Категория", "Готов"])
        writer.writerows([["PDF", "Документы", "Да"], ["PNG", "Изображения", "Да"]])


def write_archives() -> None:
    with zipfile.ZipFile(OUTPUT / "archive.zip", "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("readme.txt", "Archive preview")
        archive.writestr("folder/data.json", '{"ok": true}')
    with tarfile.open(OUTPUT / "archive.tar.gz", "w:gz") as archive:
        archive.add(OUTPUT / "hello.txt", arcname="hello.txt")
    with gzip.open(OUTPUT / "single.txt.gz", "wb") as stream:
        stream.write(b"Compressed text stream")


def write_sqlite() -> None:
    path = OUTPUT / "demo.sqlite"
    path.unlink(missing_ok=True)
    connection = sqlite3.connect(path)
    connection.execute("CREATE TABLE formats (name TEXT, category TEXT, score INTEGER)")
    connection.executemany(
        "INSERT INTO formats VALUES (?, ?, ?)",
        [("PDF", "documents", 10), ("PNG", "images", 10), ("MP4", "video", 10)],
    )
    connection.commit()
    connection.close()


def write_email_and_contacts() -> None:
    message = EmailMessage()
    message["From"] = "demo@example.test"
    message["To"] = "team@example.test"
    message["Subject"] = "Vibe Viewer demo"
    message.set_content("Это письмо отображается внутри программы.")
    (OUTPUT / "message.eml").write_bytes(message.as_bytes())
    (OUTPUT / "contact.vcf").write_text(
        "BEGIN:VCARD\nVERSION:3.0\nFN:Demo User\nEMAIL:demo@example.test\nEND:VCARD\n",
        encoding="utf-8",
    )
    (OUTPUT / "event.ics").write_text(
        "BEGIN:VCALENDAR\nVERSION:2.0\nBEGIN:VEVENT\nSUMMARY:Demo\nDTSTART:20260101T120000Z\nEND:VEVENT\nEND:VCALENDAR\n",
        encoding="utf-8",
    )


def write_wav() -> None:
    sample_rate = 8_000
    duration_samples = sample_rate // 4
    with wave.open(str(OUTPUT / "silence.wav"), "wb") as stream:
        stream.setnchannels(1)
        stream.setsampwidth(2)
        stream.setframerate(sample_rate)
        stream.writeframes(struct.pack("<h", 0) * duration_samples)


def write_extended_text_samples() -> None:
    """Create small, valid samples for text-based specialist viewers."""
    samples = {
        "captions.srt": "1\n00:00:00,000 --> 00:00:02,000\nVibe Viewer subtitle\n",
        "captions.vtt": "WEBVTT\n\n00:00:00.000 --> 00:00:02.000\nWebVTT subtitle\n",
        "captions.ass": "[Script Info]\nTitle: Demo\n[Events]\nDialogue: 0,0:00:00.00,0:00:02.00,Default,,0,0,0,,ASS subtitle\n",
        "captions.ssa": "[Script Info]\nTitle: Demo\n[Events]\nDialogue: 0,0:00:00.00,0:00:02.00,Default,,0,0,0,,SSA subtitle\n",
        "captions.sub": "{0}{50}MicroDVD subtitle\n",
        "lyrics.lrc": "[00:00.00]Vibe Viewer\n[00:02.00]Synchronized lyrics\n",
        "playlist.m3u": "#EXTM3U\n#EXTINF:2,Demo audio\nsilence.wav\n",
        "playlist.m3u8": "#EXTM3U\n#EXTINF:2,Demo audio\nsilence.wav\n",
        "playlist.pls": "[playlist]\nFile1=silence.wav\nTitle1=Demo audio\nNumberOfEntries=1\n",
        "playlist.xspf": '<?xml version="1.0"?><playlist xmlns="http://xspf.org/ns/0/"><trackList><track><title>Demo</title><location>silence.wav</location></track></trackList></playlist>',
        "album.cue": 'FILE "silence.wav" WAVE\n  TRACK 01 AUDIO\n    TITLE "Demo"\n    INDEX 01 00:00:00\n',
        "route.gpx": '<?xml version="1.0"?><gpx version="1.1"><trk><name>Demo</name><trkseg><trkpt lat="52.28" lon="104.28"/><trkpt lat="52.29" lon="104.30"/></trkseg></trk></gpx>',
        "places.kml": '<?xml version="1.0"?><kml xmlns="http://www.opengis.net/kml/2.2"><Placemark><name>Demo</name><Point><coordinates>104.28,52.28,0</coordinates></Point></Placemark></kml>',
        "geometry.gml": '<gml:Point xmlns:gml="http://www.opengis.net/gml"><gml:coordinates>104.28,52.28</gml:coordinates></gml:Point>',
        "geometry.wkt": "POINT (104.28 52.28)\nLINESTRING (104.28 52.28, 104.30 52.29)\n",
        "model.obj": "o Pyramid\nv 0 0 0\nv 1 0 0\nv 0 1 0\nv 0 0 1\nf 1 2 3\nf 1 2 4\n",
        "model.off": "OFF\n4 2 0\n0 0 0\n1 0 0\n0 1 0\n0 0 1\n3 0 1 2\n3 0 1 3\n",
        "model.ply": "ply\nformat ascii 1.0\nelement vertex 3\nproperty float x\nproperty float y\nproperty float z\nelement face 1\nproperty list uchar int vertex_indices\nend_header\n0 0 0\n1 0 0\n0 1 0\n3 0 1 2\n",
        "model.gltf": json.dumps({"asset": {"version": "2.0"}, "scenes": [{"nodes": []}], "scene": 0}),
        "settings.env": "VIEWER_MODE=demo\nLANG=ru_RU.UTF-8\n",
        "query.graphql": "query Viewer { formats { name extension } }\n",
        "message.proto": 'syntax = "proto3"; message Format { string name = 1; }\n',
        "contract.sol": "pragma solidity ^0.8.0; contract Viewer { string public name = 'Vibe'; }\n",
        "translations.po": 'msgid "Open"\nmsgstr "Открыть"\n',
        "template.pot": 'msgid "Format"\nmsgstr ""\n',
        "metadata.csvw": json.dumps({"@context": "http://www.w3.org/ns/csvw", "url": "table.csv"}),
        "notebook.ipynb": json.dumps({"cells": [{"cell_type": "markdown", "metadata": {}, "source": ["# Demo"]}], "metadata": {}, "nbformat": 4, "nbformat_minor": 5}),
    }
    for name, contents in samples.items():
        (OUTPUT / name).write_text(contents, encoding="utf-8")
    with zipfile.ZipFile(OUTPUT / "places.kmz", "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("doc.kml", samples["places.kml"])


def write_extended_binary_samples() -> None:
    """Create valid compact binary/container samples where stdlib is sufficient."""
    # Empty WebAssembly module and Java 8 class header (header inspection samples).
    (OUTPUT / "empty.wasm").write_bytes(b"\0asm\x01\0\0\0")
    (OUTPUT / "header.class").write_bytes(b"\xca\xfe\xba\xbe\0\0\0\x34\0\x01")
    # One-packet little-endian Ethernet PCAP.
    packet = bytes.fromhex("ffffffffffff0011223344550800") + bytes(46)
    pcap = struct.pack("<IHHIIII", 0xA1B2C3D4, 2, 4, 0, 0, 65535, 1)
    pcap += struct.pack("<IIII", 1_700_000_000, 0, len(packet), len(packet)) + packet
    (OUTPUT / "network.pcap").write_bytes(pcap)
    # POSIX ar archive containing a single text member; .deb uses the same container.
    payload = b"Vibe Viewer package sample\n"
    header = f"readme.txt/     {0:<12}{0:<6}{0:<6}{0o100644:<8}{len(payload):<10}`\n".encode("ascii")
    ar_data = b"!<arch>\n" + header + payload + (b"\n" if len(payload) % 2 else b"")
    (OUTPUT / "library.ar").write_bytes(ar_data)
    (OUTPUT / "package.deb").write_bytes(ar_data)
    # Minimal bencoded torrent metadata.
    (OUTPUT / "metadata.torrent").write_bytes(b"d8:announce24:https://example.invalid/4:infod4:name8:demo.txt6:lengthi42eee")
    # Minimal CPIO newc archive with one member and trailer.
    def cpio_member(name: str, contents: bytes, mode: int = 0o100644) -> bytes:
        name_bytes = name.encode() + b"\0"
        fields = (1, mode, 0, 0, 1, 0, len(contents), 0, 0, 0, 0, len(name_bytes), 0)
        header = b"070701" + b"".join(f"{value:08x}".encode() for value in fields)
        record = header + name_bytes
        record += bytes((-len(record)) % 4)
        record += contents
        return record + bytes((-len(record)) % 4)

    cpio = cpio_member("readme.txt", b"Vibe Viewer CPIO\n") + cpio_member("TRAILER!!!", b"")
    (OUTPUT / "archive.cpio").write_bytes(cpio)
    # Minimal CAB header suitable for safe header inspection.
    cab = bytearray(36)
    cab[:4] = b"MSCF"
    struct.pack_into("<I", cab, 8, 36)
    struct.pack_into("<I", cab, 16, 36)
    cab[24:26] = bytes((3, 1))
    (OUTPUT / "archive.cab").write_bytes(cab)


def write_optional_samples() -> None:
    try:
        from PIL import Image, ImageDraw

        image = Image.new("RGB", (640, 360), "#dbeafe")
        draw = ImageDraw.Draw(image)
        draw.rectangle((40, 40, 600, 320), outline="#2563eb", width=6)
        draw.text((80, 155), "Vibe Viewer image", fill="#172033")
        image.save(OUTPUT / "image.png")
        image.save(OUTPUT / "image.webp")
    except ImportError:
        pass

    try:
        import trimesh

        mesh = trimesh.creation.box()
        for extension in ("stl", "glb", "3mf"):
            mesh.export(OUTPUT / f"model.{extension}")
    except (ImportError, ValueError):
        pass

    try:
        from astropy.io import fits
        import numpy as np

        fits.PrimaryHDU(np.arange(100, dtype="int16").reshape(10, 10)).writeto(OUTPUT / "image.fits", overwrite=True)
        for extension in ("fts", "fit"):
            (OUTPUT / f"image.{extension}").write_bytes((OUTPUT / "image.fits").read_bytes())
    except ImportError:
        pass

    try:
        import netCDF4

        with netCDF4.Dataset(OUTPUT / "climate.nc", "w") as dataset:
            dataset.createDimension("time", 3)
            variable = dataset.createVariable("temperature", "f4", ("time",))
            variable[:] = [18.5, 19.0, 20.25]
        (OUTPUT / "climate.cdf").write_bytes((OUTPUT / "climate.nc").read_bytes())
    except ImportError:
        pass

    try:
        import scipy.io
        import numpy as np

        scipy.io.savemat(OUTPUT / "matrix.mat", {"demo": np.arange(12).reshape(3, 4)})
    except ImportError:
        pass

    try:
        import shapefile

        with shapefile.Writer(str(OUTPUT / "points.shp")) as writer:
            writer.field("name", "C")
            writer.point(104.28, 52.28)
            writer.record("Irkutsk")
    except ImportError:
        pass

    try:
        import zstandard
        import lz4.frame

        data = b"Modern compressed stream sample"
        (OUTPUT / "stream.zst").write_bytes(zstandard.ZstdCompressor().compress(data))
        (OUTPUT / "stream.lz4").write_bytes(lz4.frame.compress(data))
    except ImportError:
        pass

    try:
        import pandas as pd
        import pyreadstat

        frame = pd.DataFrame({"format": ["PDF", "PNG"], "score": [10, 9]})
        pyreadstat.write_sav(frame, OUTPUT / "statistics.sav", file_label="Vibe Viewer demo")
    except ImportError:
        pass

    try:
        from dbfread import DBF  # noqa: F401 - verifies that the reader is installed
        import shapefile

        with shapefile.Writer(dbf=str(OUTPUT / "table.dbf")) as writer:
            writer.field("FORMAT", "C")
            writer.field("SCORE", "N")
            writer.record("PDF", 10)
            writer.record("PNG", 9)
    except (ImportError, TypeError):
        pass

    try:
        from fastavro import writer

        schema = {"type": "record", "name": "Format", "fields": [{"name": "name", "type": "string"}, {"name": "score", "type": "int"}]}
        with (OUTPUT / "table.avro").open("wb") as stream:
            writer(stream, schema, [{"name": "PDF", "score": 10}, {"name": "PNG", "score": 9}])
    except ImportError:
        pass

    try:
        import pyarrow as pa
        import pyarrow.orc as orc

        orc.write_table(pa.table({"name": ["PDF", "PNG"], "score": [10, 9]}), OUTPUT / "table.orc")
    except ImportError:
        pass

    try:
        import pycdlib

        iso = pycdlib.PyCdlib()
        iso.new(interchange_level=3)
        iso.add_fp(__import__("io").BytesIO(b"Vibe Viewer ISO"), len(b"Vibe Viewer ISO"), iso_path="/README.TXT;1")
        iso.write(str(OUTPUT / "disc.iso"))
        iso.close()
    except ImportError:
        pass

    try:
        import pymupdf

        document = pymupdf.open()
        page = document.new_page()
        page.insert_text((72, 100), "Vibe Viewer PDF demo", fontsize=24)
        document.save(OUTPUT / "document.pdf")
        document.close()
    except ImportError:
        pass

    try:
        from docx import Document

        document = Document()
        document.add_heading("Vibe Viewer DOCX", 0)
        document.add_paragraph("Документ открыт без Microsoft Word и LibreOffice.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Формат"
        table.cell(0, 1).text = "Статус"
        table.cell(1, 0).text = "DOCX"
        table.cell(1, 1).text = "Работает"
        document.save(OUTPUT / "document.docx")
    except ImportError:
        pass

    try:
        from openpyxl import Workbook

        book = Workbook()
        sheet = book.active
        sheet.title = "Форматы"
        sheet.append(["Формат", "Категория", "Готов"])
        sheet.append(["XLSX", "Таблицы", True])
        book.save(OUTPUT / "workbook.xlsx")
    except ImportError:
        pass

    try:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Vibe Viewer PPTX"
        slide.placeholders[1].text = "Предпросмотр текста и изображений слайда"
        presentation.save(OUTPUT / "slides.pptx")
    except ImportError:
        pass

    try:
        import numpy as np

        np.save(OUTPUT / "array.npy", np.arange(24).reshape(4, 6))
        np.savez(OUTPUT / "arrays.npz", first=np.arange(5), second=np.eye(3))
    except ImportError:
        pass


def main() -> int:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    write_text_samples()
    write_archives()
    write_sqlite()
    write_email_and_contacts()
    write_wav()
    write_extended_text_samples()
    write_extended_binary_samples()
    write_optional_samples()
    print(f"Generated demo files in {OUTPUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
